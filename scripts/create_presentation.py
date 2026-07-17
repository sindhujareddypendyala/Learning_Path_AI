import sys
import os
import subprocess

# Installs python-pptx if not present
try:
    import pptx
except ImportError:
    print("Installing python-pptx library...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom color palette (Matching LearnPath-AI aesthetics)
    BG_COLOR = RGBColor(15, 23, 42)       # Dark Navy (Slate-900)
    PRIMARY_COLOR = RGBColor(129, 140, 248) # Indigo (Light)
    ACCENT_COLOR = RGBColor(34, 211, 238)  # Cyan
    TEXT_COLOR = RGBColor(241, 245, 249)    # Off-white
    MUTED_COLOR = RGBColor(148, 163, 184)   # Gray
    
    # Helper to set solid background color
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper to create slides with title and content
    def add_slide(title_text):
        blank_slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        set_background(slide)
        
        # Add Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Outfit"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        return slide

    # ================= SLIDE 1: Title Slide =================
    slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(slide_layout)
    set_background(slide1)
    
    # Title & Subtitle in a single box to avoid overlap
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "LearnPath-AI"
    p1.font.name = "Outfit"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_COLOR
    p1.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Enterprise-Grade Multi-Agent Personalized Learning Platform"
    p2.font.name = "Inter"
    p2.font.size = Pt(24)
    p2.font.color.rgb = PRIMARY_COLOR
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "Powered by FastAPI, React & LangChain Multi-Agent Architecture"
    p3.font.name = "Inter"
    p3.font.size = Pt(16)
    p3.font.color.rgb = MUTED_COLOR
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(24)

    # ================= SLIDE 2: Problem & Solution =================
    slide2 = add_slide("The Problem & Our Solution")
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "The Problem: Traditional E-Learning is Broken"
    p.font.name = "Inter"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    bullet1 = tf.add_paragraph()
    bullet1.text = "• One-size-fits-all: Static courses ignore individual goals and existing skill levels."
    bullet1.font.name = "Inter"
    bullet1.font.size = Pt(18)
    bullet1.font.color.rgb = TEXT_COLOR
    bullet1.space_before = Pt(8)
    
    bullet2 = tf.add_paragraph()
    bullet2.text = "• Lack of hands-on practice: Courses focus on passive video consumption rather than building."
    bullet2.font.name = "Inter"
    bullet2.font.size = Pt(18)
    bullet2.font.color.rgb = TEXT_COLOR
    bullet2.space_before = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = "Our Solution: LearnPath-AI"
    p2.font.name = "Inter"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY_COLOR
    p2.space_before = Pt(24)
    
    bullet3 = tf.add_paragraph()
    bullet3.text = "• Personalized Pathing: Generate whole customized roadmaps from a single prompt."
    bullet3.font.name = "Inter"
    bullet3.font.size = Pt(18)
    bullet3.font.color.rgb = TEXT_COLOR
    bullet3.space_before = Pt(8)
    
    bullet4 = tf.add_paragraph()
    bullet4.text = "• Interactive Learning Loops: Specialized agents adapt modules, generate practice labs, and quiz users dynamically."
    bullet4.font.name = "Inter"
    bullet4.font.size = Pt(18)
    bullet4.font.color.rgb = TEXT_COLOR
    bullet4.space_before = Pt(6)

    # ================= SLIDE 3: Architecture =================
    slide3 = add_slide("Technical Architecture Overview")
    content_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Built on clean architecture, SOLID principles, and serverless scalability."
    p.font.name = "Inter"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    
    h1 = tf.add_paragraph()
    h1.text = "1. Unified Monorepo Workspace"
    h1.font.name = "Inter"
    h1.font.size = Pt(20)
    h1.font.bold = True
    h1.font.color.rgb = PRIMARY_COLOR
    h1.space_before = Pt(18)
    
    b1 = tf.add_paragraph()
    b1.text = "React SPA frontend served via Vite + FastAPI REST backend on python-uvicorn."
    b1.font.name = "Inter"
    b1.font.size = Pt(16)
    b1.font.color.rgb = MUTED_COLOR
    b1.space_before = Pt(4)
    
    h2 = tf.add_paragraph()
    h2.text = "2. Agentic Workflow Layer"
    h2.font.name = "Inter"
    h2.font.size = Pt(20)
    h2.font.bold = True
    h2.font.color.rgb = PRIMARY_COLOR
    h2.space_before = Pt(18)
    
    b2 = tf.add_paragraph()
    b2.text = "LangChain orchestrator directing task queues to LLM reasoning models (Llama 3 via Groq API)."
    b2.font.name = "Inter"
    b2.font.size = Pt(16)
    b2.font.color.rgb = MUTED_COLOR
    b2.space_before = Pt(4)
    
    h3 = tf.add_paragraph()
    h3.text = "3. RAG Storage Layer"
    h3.font.name = "Inter"
    h3.font.size = Pt(20)
    h3.font.bold = True
    h3.font.color.rgb = PRIMARY_COLOR
    h3.space_before = Pt(18)
    
    b3 = tf.add_paragraph()
    b3.text = "ChromaDB vector store with sentence-transformers for embedding course contexts and pdf documents."
    b3.font.name = "Inter"
    b3.font.size = Pt(16)
    b3.font.color.rgb = MUTED_COLOR
    b3.space_before = Pt(4)

    # ================= SLIDE 4: Agent Team =================
    slide4 = add_slide("The Specialized Multi-Agent Team")
    content_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Collaborating agents performing specialized tasks (coordinating via shared memory state):"
    p.font.name = "Inter"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    
    agents = [
        ("Curriculum Agent", "Analyzes learning goals to formulate modules and log sequences."),
        ("Content Agent", "Generates detailed lessons and Tip guides for code mental models."),
        ("Assessment Agent", "Constructs and evaluates interactive checkpoints and quizzes."),
        ("Projects Agent", "Designs relevant portfolio sprint tasks and review rubrics."),
        ("Interview Agent", "Assembles technical mock interview prep questions."),
        ("Analytics Agent", "Analyzes completing metrics and updates learning progress graphs.")
    ]
    
    for name, desc in agents:
        ap = tf.add_paragraph()
        ap.text = f"• {name}: {desc}"
        ap.font.name = "Inter"
        ap.font.size = Pt(16)
        ap.font.color.rgb = MUTED_COLOR
        ap.space_before = Pt(10)

    # ================= SLIDE 5: Core Features =================
    slide5 = add_slide("Key Features & Capabilities")
    content_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    features = [
        ("Instant Course Generation", "A single prompt produces structured 4-module roadmaps detailing difficulty levels, completion times, and objectives."),
        ("Interactive Quizzes & Confetti", "Checks user understanding and celebrates high scores with visual animations using canvas-confetti."),
        ("Portfolio sprint labs", "Assigns practical projects, defining criteria and tracking skills (e.g. Embeddings, APIs, python classification)."),
        ("Active Interview Simulator", "Enables mock technical interview panels with hidden answer reveals and structural guidance (context, decision, tradeoff)."),
        ("Sleek Premium UI", "Built with React, TailwindCSS, glassmorphic styles, Outfit/Inter typography, and smooth framer-motion micro-animations.")
    ]
    
    for title, desc in features:
        h = tf.add_paragraph()
        h.text = f"• {title}"
        h.font.name = "Inter"
        h.font.size = Pt(18)
        h.font.bold = True
        h.font.color.rgb = PRIMARY_COLOR
        h.space_before = Pt(12)
        
        b = tf.add_paragraph()
        b.text = desc
        b.font.name = "Inter"
        b.font.size = Pt(15)
        b.font.color.rgb = TEXT_COLOR
        b.space_before = Pt(2)

    # ================= SLIDE 6: Deployment =================
    slide6 = add_slide("Production Deployment Strategies")
    content_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Highly modular setup allows deploying easily on any cloud provider:"
    p.font.name = "Inter"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    
    h1 = tf.add_paragraph()
    h1.text = "1. Containerized VPS Deployment (Docker Compose)"
    h1.font.name = "Inter"
    h1.font.size = Pt(18)
    h1.font.bold = True
    h1.font.color.rgb = PRIMARY_COLOR
    h1.space_before = Pt(16)
    
    b1 = tf.add_paragraph()
    b1.text = "Uses Dockerfile.frontend (Nginx multi-stage build) & Dockerfile.backend (FastAPI/LangChain setup) orchestrated with local persistent database volumes."
    b1.font.name = "Inter"
    b1.font.size = Pt(15)
    b1.font.color.rgb = MUTED_COLOR
    b1.space_before = Pt(4)
    
    h2 = tf.add_paragraph()
    h2.text = "2. Serverless / Managed Cloud Deployment"
    h2.font.name = "Inter"
    h2.font.size = Pt(18)
    h2.font.bold = True
    h2.font.color.rgb = PRIMARY_COLOR
    h2.space_before = Pt(16)
    
    b2 = tf.add_paragraph()
    b2.text = "Deploy the backend on Render / Hugging Face Spaces (CPU basic Docker Space) and the frontend on Vercel (using vercel.json rewrite API proxy rules)."
    b2.font.name = "Inter"
    b2.font.size = Pt(15)
    b2.font.color.rgb = MUTED_COLOR
    b2.space_before = Pt(4)

    # ================= SLIDE 7: Summary =================
    slide7 = add_slide("Conclusion & Future Scope")
    content_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullet1 = tf.add_paragraph()
    bullet1.text = "• Summary: LearnPath-AI represents a premium, production-ready implementation of a multi-agent learning management platform."
    bullet1.font.name = "Inter"
    bullet1.font.size = Pt(18)
    bullet1.font.color.rgb = TEXT_COLOR
    
    bullet2 = tf.add_paragraph()
    bullet2.text = "• Future Scope: Integrations with third-party LMS providers via LTI standards."
    bullet2.font.name = "Inter"
    bullet2.font.size = Pt(18)
    bullet2.font.color.rgb = TEXT_COLOR
    bullet2.space_before = Pt(16)
    
    bullet3 = tf.add_paragraph()
    bullet3.text = "• Automated Code Review: Incorporating sandbox code execution to evaluate user-submitted projects dynamically."
    bullet3.font.name = "Inter"
    bullet3.font.size = Pt(18)
    bullet3.font.color.rgb = TEXT_COLOR
    bullet3.space_before = Pt(16)
    
    bullet4 = tf.add_paragraph()
    bullet4.text = "• Multi-modal Content: Generating dynamic audio/video summaries alongside text-based lessons."
    bullet4.font.name = "Inter"
    bullet4.font.size = Pt(18)
    bullet4.font.color.rgb = TEXT_COLOR
    bullet4.space_before = Pt(16)

    # Save presentation
    output_path = "c:\\Users\\sindh\\OneDrive\\Desktop\\LearnPath_AI_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_presentation()
